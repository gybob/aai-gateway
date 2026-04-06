import sys
from pptx import Presentation
import copy

pptx_file = sys.argv[1]
out_dir = sys.argv[2]

prs = Presentation(pptx_file)
slide_width = prs.slide_width
slide_height = prs.slide_height

for i, slide in enumerate(prs.slides):
    new_prs = Presentation()
    new_prs.slide_width = slide_width
    new_prs.slide_height = slide_height

    blank_layout = new_prs.slide_layouts[6]
    new_slide = new_prs.slides.add_slide(blank_layout)

    if slide.background.fill.type is not None:
        try:
            new_slide.background.fill.solid()
            new_slide.background.fill.fore_color.rgb = slide.background.fill.fore_color.rgb
        except:
            pass

    for shape in slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    out_path = f'{out_dir}/slide_{i+1}.pptx'
    new_prs.save(out_path)
    print(f'Split slide {i+1}')
